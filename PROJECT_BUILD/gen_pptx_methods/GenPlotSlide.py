from collections.abc import Container
from numpy.lib.npyio import loadtxt
from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
import logging.config

logging.config.fileConfig("../debug/logging.conf")
logger = logging.getLogger("GenPlotSlide")


class GenPlotSlide:
    def __init__(self, title, dat_path, pageNum, filePath):
        self.TITLE = title
        self.DAT_PATH = dat_path
        self.PAGE_NUM = int(pageNum)
        self.FILE_PATH = filePath

    def generate(self):
        ppt = Presentation(self.FILE_PATH)

        plot_slide_register = ppt.slide_layouts[5]  # plot slide

        plot_slide = ppt.slides.add_slide(plot_slide_register)

        title = plot_slide.shapes.title
        title.text = self.TITLE
        logger.debug("generated title")

        # retreive dat file
        logger.debug(f"retreiving {self.DAT_PATH} file")
        x, y = loadtxt(self.DAT_PATH, unpack=True, usecols=[0, 1])
        xydict = dict(zip(x, y))

        # build graph
        chart_data = XyChartData()
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        series_1 = chart_data.add_series('Model 1')
        for i in xydict:
            series_1.add_data_point(i, xydict[i])

        chart = plot_slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
        ).chart
        logger.debug("plot generation complete!")

        # add page number
        logger.debug("adding page number")
        tx_box = plot_slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
        tf = tx_box.text_frame
        pn = tf.add_paragraph()
        pn.text = f"{self.PAGE_NUM}"
        pn.font.size = Pt(15)

        ppt.save(self.FILE_PATH)
        logger.debug("generation complete!")
