from collections.abc import Container
from pptx import Presentation
from pptx.util import Inches, Pt
# import logging
import logging.config

logging.config.fileConfig("../debug/logging.conf")
logger = logging.getLogger("GenTitleSlide")


class GenTitleSlide:
    def __init__(self, title, subtitle, pageNum, filePath):
        self.TITLE = title
        self.SUBTITLE = subtitle
        self.PAGE_NUM = int(pageNum)
        self.FILE_PATH = filePath

    def generate(self):
        logger.debug(f"|title: {self.TITLE} |subtitle: {self.SUBTITLE} "
                     f"|path: {self.FILE_PATH} |page number: {self.PAGE_NUM}")

        ppt = Presentation(self.FILE_PATH)

        title_slide_register = ppt.slide_layouts[0]  # title page

        title_slide = ppt.slides.add_slide(title_slide_register)

        title = title_slide.shapes.title

        subtitle = title_slide.placeholders[1]

        title.text = self.TITLE
        subtitle.text = self.SUBTITLE

        # add page number
        tx_box = title_slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
        tf = tx_box.text_frame
        pn = tf.add_paragraph()
        pn.text = f"{self.PAGE_NUM}"
        pn.font.size = Pt(15)

        ppt.save(self.FILE_PATH)
        logger.debug("generation complete!")
