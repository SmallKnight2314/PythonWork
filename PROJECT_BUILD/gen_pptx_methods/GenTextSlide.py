from collections.abc import Container
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
# import logging
import logging.config

logging.config.fileConfig("../debug/logging.conf")
logger = logging.getLogger("GenTextSlide")


class GenTextSlide:
    def __init__(self, title, paragraphs, pageNum, filePath):
        self.TITLE = title
        self.PARAGRAPHS = paragraphs
        self.PAGE_NUM = int(pageNum)
        self.FILE_PATH = filePath

    def generate(self):
        logger.debug(f"|title: {self.TITLE} "
                     f"|path: {self.FILE_PATH} |page number: {self.PAGE_NUM}")

        ppt = Presentation(self.FILE_PATH)

        text_slide_register = ppt.slide_layouts[5]

        text_slide = ppt.slides.add_slide(text_slide_register)
        title = text_slide.shapes.title
        title.text = self.TITLE

        left, top, width, height = Inches(0.5), Inches(2), Inches(9), Inches(4.5)

        paragraph_cont = text_slide.shapes.add_textbox(left, top, width, height)
        tf = paragraph_cont.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True
        tf.margin_bottom = Inches(0.08)
        tf.margin_left = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        logger.debug("convert dict into list of text")
        paragraphs = self.PARAGRAPHS
        paragraph_strs = []
        for para in paragraphs:
            paragraph_strs.append(para["text"])
        logger.debug(f"full paragraphs_strs : \n{paragraph_strs}")

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        logger.debug("adding paragraphs")
        for para_str in paragraph_strs[0:]:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "\t" + para_str + "\n"
            p.alignment = PP_ALIGN.LEFT
            p.level = 0

        # add page number
        logger.debug("adding page number")
        tx_box = text_slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
        tf = tx_box.text_frame
        pn = tf.add_paragraph()
        pn.text = f"{self.PAGE_NUM}"
        pn.font.size = Pt(15)
        logger.debug("added page number to slide")

        ppt.save(self.FILE_PATH)
        logger.debug("generation complete!")
