from collections.abc import Container
from pptx import Presentation
from pptx.util import Inches, Pt
import logging.config

logging.config.fileConfig("../debug/logging.conf")
logger = logging.getLogger("GenListSlide")


class GenListSlide:
    def __init__(self, title, lists, pageNum, filePath):
        self.TITLE = title
        self.LISTS = lists  # special list, is dictionary/list hybrid
        self.PAGE_NUM = int(pageNum)
        self.FILE_PATH = filePath

    def generate(self):
        ppt = Presentation(self.FILE_PATH)

        list_slide_register = ppt.slide_layouts[1]

        list_slide = ppt.slides.add_slide(list_slide_register)

        title = list_slide.shapes.title
        title.text = self.TITLE

        bullet_point_box = list_slide.shapes
        bullet_points_container = bullet_point_box.placeholders[1]  # level 0
        bullet_points_container.text = "List items:"
        logger.debug("begin list generation")
        # adding list items and its children
        for listItems in self.LISTS:
            logger.debug("check for sublists")
            if listItems["subListItems"] == False:
                logger.debug("generating lonely list item")
                bullet_points_lvl0 = bullet_points_container.text_frame.add_paragraph()
                bullet_points_lvl0.text = listItems["listItem"]
                bullet_points_lvl0.level = 1
                logger.debug("generated lonely list item")
            else:
                logger.debug("generating complex list item")
                bullet_points_lvl0 = bullet_points_container.text_frame.add_paragraph()  # level 1
                bullet_points_lvl0.text = listItems["listItem"]
                bullet_points_lvl0.level = 1
                logger.debug("generation complete, moving to sublist generation")
                for subItems in listItems["subList"]:
                    logger.debug("generate sublist item")
                    bullet_points_lvl1 = bullet_points_container.text_frame.add_paragraph()  # level 2
                    bullet_points_lvl1.text = subItems["subListItem"]
                    bullet_points_lvl1.level = 2
                    logger.debug("sublist complete")
            logger.debug("cycle repeat")
        logger.debug("cycle complete!")

        # add page number
        logger.debug("adding page number")
        tx_box = list_slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
        tf = tx_box.text_frame
        pn = tf.add_paragraph()
        pn.text = f"{self.PAGE_NUM}"
        pn.font.size = Pt(15)
        logger.debug("added page number to slide")

        ppt.save(self.FILE_PATH)
        logger.debug("generation complete!")
