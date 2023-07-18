from collections.abc import Container
from pptx import Presentation
from pptx.util import Inches, Pt
import logging.config

logging.config.fileConfig("../debug/logging.conf")
logger = logging.getLogger("GenPictureSlide")


class GenPictureSlide:
    def __init__(self, title, img_url, pageNum, caption, filePath):
        self.TITLE = title
        self.IMG_URL = img_url
        self.PAGE_NUM = int(pageNum)
        self.FILE_PATH = filePath
        self.CAPTION = caption

    def generate(self):
        logger.debug("begin picture slide generation")
        ppt = Presentation(self.FILE_PATH)

        picture_slide_register = ppt.slide_layouts[8]

        picture_slide = ppt.slides.add_slide(picture_slide_register)
        title = picture_slide.shapes.title
        _picture = picture_slide.placeholders[1]
        content = picture_slide.placeholders[2]

        title.text = self.TITLE
        logger.debug("added title")
        _picture.insert_picture(self.IMG_URL)
        logger.debug("added picture")
        content.text = self.CAPTION
        logger.debug("added caption")

        # add page number
        logger.debug("adding page number")
        tx_box = picture_slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
        tf = tx_box.text_frame
        pn = tf.add_paragraph()
        pn.text = f"{self.PAGE_NUM}"
        pn.font.size = Pt(15)
        logger.debug("added page number to slide")

        ppt.save(self.FILE_PATH)
        logger.debug("generation complete!")
